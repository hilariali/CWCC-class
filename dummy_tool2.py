# dummy_tool2.py

import streamlit as st
import openai
from io import BytesIO
from docx import Document
from pptx import Presentation
import PyPDF2

# Config
MAX_FILES = 10
MAX_CHAR_LIMIT = 400_000  # Safe buffer for 130k token limit

def extract_text_from_file(file, file_type):
    try:
        if file_type == "txt":
            return file.read().decode("utf-8")

        elif file_type == "pdf":
            reader = PyPDF2.PdfReader(file)
            return "\n".join([page.extract_text() or "" for page in reader.pages])

        elif file_type in ["doc", "docx"]:
            doc = Document(file)
            return "\n".join([para.text for para in doc.paragraphs])

        elif file_type in ["ppt", "pptx"]:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
    except Exception:
        return ""
    return ""

def run():
    st.title("📁 Chat with Your Files")
    st.caption("Upload documents and chat with AI about their contents.")

    client = openai.OpenAI(
        api_key=st.secrets["OPENAI_API_KEY"],
        base_url=st.secrets.get("OPENAI_BASE_URL"),
    )

    # Session state for context and chat
    if "file_context" not in st.session_state:
        st.session_state.file_context = ""
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = [
            {"role": "assistant", "content": "Hi! Upload your documents and ask me anything about them."}
        ]

    uploaded_files = st.file_uploader(
        "Upload up to 10 files (TXT, PDF, DOCX, PPTX)", 
        type=["txt", "pdf", "doc", "docx", "ppt", "pptx"],
        accept_multiple_files=True
    )

    if uploaded_files:
        if len(uploaded_files) > MAX_FILES:
            st.warning("⚠️ You can upload a maximum of 10 files.")
            return

        combined_text = ""
        total_chars = 0

        for file in uploaded_files:
            ext = file.name.split(".")[-1].lower()
            content = extract_text_from_file(file, ext)
            if not content:
                st.error(f"❌ Could not extract text from: {file.name}")
                continue

            content_len = len(content)
            if total_chars + content_len > MAX_CHAR_LIMIT:
                st.warning(f"⚠️ Skipping {file.name}: it would exceed the token limit.")
                continue

            combined_text += f"\n\n--- {file.name} ---\n{content}"
            total_chars += content_len
            st.success(f"✓ Extracted: {file.name} ({content_len:,} characters)")

        st.session_state.file_context = combined_text
        st.info(f"🧮 Total characters extracted: {total_chars:,} (~{total_chars // 4:,} tokens)")

        if combined_text:
            st.subheader("📄 File Text Preview")
            st.text_area("Combined file text (first 5000 characters shown):", value=combined_text[:5000], height=300)

    # Model choice
    model_options = [
        "meta-llama/Llama-3.3-70B-Instruct"
        #"DeepSeek-R1-Distill-Qwen-32B",
        #"DeepSeek-V3-1",
        #"gpt-oss-120b",
        #"Meta-Llama-3-1-8B-Instruct-FP8",
        #"Meta-Llama-3-3-70B-Instruct",
        #"Meta-Llama-4-Maverick-17B-128E-Instruct-FP8",
        #"Qwen3-235B-A22B-Instruct-2507-FP8",
    ]
    selected_model = st.selectbox("Choose an AI model:", model_options)

    st.divider()

    # Display chat history
    for msg in st.session_state.chat_history:
        st.chat_message(msg["role"]).write(msg["content"])

    user_input = st.chat_input("Ask something about the uploaded files…")
    if user_input and not st.session_state.file_context:
        st.warning("Please upload and process files before chatting.")
    elif user_input:
        st.session_state.chat_history.append({"role": "user", "content": user_input})
        st.chat_message("user").write(user_input)

        with st.spinner("Thinking..."):
            try:
                system_prompt = (
                    f"You are an assistant that answers questions based on the following uploaded documents:\n\n"
                    f"{st.session_state.file_context[:MAX_CHAR_LIMIT]}\n\n"
                    f"Now answer based on this information."
                )
                full_history = [{"role": "system", "content": system_prompt}] + st.session_state.chat_history

                response = client.chat.completions.create(
                    model="meta-llama/Llama-3.3-70B-Instruct",
                    messages=full_history,
                )
                reply = response.choices[0].message.content
                st.session_state.chat_history.append({"role": "assistant", "content": reply})
                st.chat_message("assistant").write(reply)
            except Exception as e:
                st.error(f"❌ Error: {e}")
